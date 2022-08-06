import magic
import olefile                  # 한글
from pptx import Presentation   # ppt 
from docx import Document       # 워드
import openpyxl as oxl          # 엑셀
import pdfplumber               # pdf


class FileReader():
    """
        문서 파일 읽기 클래스
    """
    def __init__(self):
        
        pass


        
    def read_data(self, file_path):
        file_type = self.__get_file_type(file_path)
        
        if file_type is None:
            return False
        if 'Hangul' in file_type: 
            result = self.__read_hwp(file_path)
            return result
        elif 'Word' in file_type:
            result = self.__read_word(file_path)
            return result
        elif 'PDF' in file_type:
            result = self.__read_pdf(file_path)
            return result
        elif 'excel' in file_type or 'csv' in file_type or 'Zip archive data' in file_type:
            result = self.__read_excel(file_path)
            return result
        elif 'PowerPoint' in file_type:
            result = self.__read_ppt(file_path)
            return result
        else:
            result = self.__read_text(file_path)
            return result
            
            
    def __get_file_type(self, file_path):
        try:
            file_type = magic.from_file(file_path)
            return file_type
        except:
            
            file_info = {file_path:file_type}
            return file_info

    def __read_hwp(self, file_path):
        file = olefile.OleFileIO(file_path)
        
        encoded_data = file.openstream('PrvText').read()
        decoded_data = encoded_data.decode('utf-16')

        # img = file.openstream('PrvImage').read()
        data = ''.join(map(str, decoded_data))
        return data


    def __read_word(self, file_path):
        data = []
        doc = Document(file_path)
        text_data = [paragraph.text for paragraph in doc.paragraphs]

        tables = doc.tables
        tables_data = []
        table_data = [cell.text for table in tables for row in table.rows for cell in row.cells]
        
        for text in text_data:
            data.append(text)
        
        # data.append(text_data for txt_data in text_data)
        for tb_data in table_data:
            data.append(tb_data)
        # data.append(tb_data for tb_data in table_data)
        
        data = '\n'.join(map(str, data))

        return data
        

    def __read_pdf(self, file_path):
        
        pdf = pdfplumber.open(file_path)
        pages = pdf.pages
        data =[]

        for page in pages:
            data.append((lambda page: page.extract_text())(page))

        data = ''.join(map(str, data))
        
        return data

    # def __read_excel(self, file_path):
    #     wb = oxl.load_workbook(file_path, data_only=True)
    #     ws = wb.active

    #     all_values = [[cell.value for cell in row]for row in ws.rows]   
    #     data = sum(all_values, [])                                # 1차원 리스트로 변환
        
    #     data = ''.join(map(str, data))
    #     # print(all_values)
    #     return data

    
    # def __read_ppt(self, file_path):
    #     data = []
    #     prs = Presentation(file_path)


    #     for slide in prs.slides:
    #         for shape in slide.shapes:
    #             if not shape.has_text_frame:
    #                 continue
    #             for paragraph in shape.text_frame.paragraphs:
    #                 data.append(paragraph.text)
        
    #     data = ''.join(map(str, data))
    #     return data


    def __read_text(self, file_path):
        data =[]
        with open(file_path, "r", encoding="utf8") as f:
            lines = f.readlines()
            for line in lines:
                data.append(line)
        data = ''.join(map(str, data))
        return data

if __name__ == '__main__':

    file_path = "./explorer_lib/test/test.txt"

    reader = FileReader()
    data = reader.read_data(file_path)
    print(data)